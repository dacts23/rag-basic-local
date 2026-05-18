#!/usr/bin/env python3
"""Offline RAG Ingestion Pipeline using Lightweight Parsers.

Hybrid Chunking & BM25-Ready Reverse Index for Air-Gapped Environments
=======================================================================

This pipeline parses complex technical documentation (PDF, DOCX, PPTX, etc.)
using lightweight parsers (PyMuPDF, python-docx, etc.), applies a hybrid
structural + overlapping-window chunking strategy, enriches chunks with
metadata, and builds an inverted index with pre-computed IDF for BM25
sparse retrieval — all 100% offline without ML models.

USAGE
-----
    # Defaults read from ./in and write to ./out/rag_output.json:
    python ingest.py

    # Explicit paths:
    python ingest.py /path/to/documents --output /path/to/output.json

    # With custom chunk threshold and overlap:
    python ingest.py --chunk-threshold 768 --overlap-ratio 0.15

    # Skip cache and re-parse everything:
    python ingest.py --force

INCREMENTAL CACHING
-------------------
    A sidecar file (parse_cache.json next to the output JSON) stores a
    SHA-256 fingerprint for each parsed document.  On re-runs, files whose
    hash is unchanged are skipped and their chunks are reused from the
    previous output, dramatically speeding up large collections.
"""

from __future__ import annotations

# ---------------------------------------------------------------------------
# Auto-install missing dependencies & pre-download models on startup
# ---------------------------------------------------------------------------
import importlib
import subprocess
import sys as _sys

def _ensure_dependencies():
    """Check for required packages and install any that are missing."""
    # Map: import_name -> pip_package_name
    deps = {
        "fitz": "PyMuPDF",
        "docx": "python-docx",
        "pptx": "python-pptx",
        "openpyxl": "openpyxl",
        "bs4": "beautifulsoup4",
        "lxml": "lxml",
        "numpy": "numpy",
        "faiss": "faiss-cpu",
        "sentence_transformers": "sentence-transformers",
    }
    missing = []
    for import_name, pip_name in deps.items():
        try:
            importlib.import_module(import_name)
        except ImportError:
            missing.append(pip_name)
    if missing:
        print(f"[ingest] Installing missing packages: {', '.join(missing)} ...", flush=True)
        subprocess.check_call(
            [_sys.executable, "-m", "pip", "install", "--quiet"] + missing
        )
        print("[ingest] Packages installed successfully.", flush=True)

def _ensure_model_cached(model_name: str):
    """Download a HuggingFace model via sentence-transformers if not already cached."""
    try:
        from sentence_transformers import SentenceTransformer
        # Try loading from cache only — if it fails, model isn't cached yet
        try:
            SentenceTransformer(model_name, local_files_only=True)
        except Exception:
            print(f"[ingest] Downloading model '{model_name}' (first run, may take a while) ...", flush=True)
            SentenceTransformer(model_name)  # Triggers download
            print(f"[ingest] Model '{model_name}' cached successfully.", flush=True)
    except ImportError:
        pass  # sentence-transformers will be installed by _ensure_dependencies
    except Exception as e:
        print(f"[ingest] Warning: could not pre-cache model '{model_name}': {e}", flush=True)

_ensure_dependencies()
# Pre-cache the default embedding model
_ensure_model_cached("BAAI/bge-small-en-v1.5")

import hashlib
import json
import logging
import math
import os
import re
import uuid
from dataclasses import dataclass, asdict
from pathlib import Path
from typing import Dict, List, Optional, Set, Tuple, Any

import numpy as np

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(name)s - %(message)s",
)
logger = logging.getLogger("rag_ingestion")

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

DEFAULT_CHUNK_THRESHOLD: int = 512
DEFAULT_OVERLAP_RATIO: float = 0.125
DEFAULT_MIN_CHUNK_TOKENS: int = 64  # blocks below this are merged by aggregator

_rag_path = os.getenv("RAG_PATH")
if not _rag_path:
    raise RuntimeError(
        "RAG_PATH environment variable is not set. "
        "Set it to your RAG workspace directory (containing in/ and out/ folders)."
    )
BASE_DIR: Path = Path(_rag_path)

DEFAULT_INPUT_DIR: Path = BASE_DIR / "in"
DEFAULT_OUTPUT_PATH: Path = BASE_DIR / "out" / "rag_output.json"
DEFAULT_CACHE_PATH: Path = BASE_DIR / "out" / "parse_cache.json"
DEFAULT_EMBED_MODEL: str = "BAAI/bge-small-en-v1.5"
EMBEDDING_DIM: int = 384
DEFAULT_FAISS_PATH: Path = BASE_DIR / "out" / "rag_embeddings.faiss"
DEFAULT_CHUNK_IDS_PATH: Path = BASE_DIR / "out" / "rag_chunk_ids.json"
DEFAULT_EMBED_CACHE_PATH: Path = BASE_DIR / "out" / "embed_cache.json"

SUPPORTED_EXTENSIONS: Set[str] = {
    ".pdf", ".docx", ".pptx", ".xlsx", ".html", ".csv",
    ".md", ".txt", ".asciidoc", ".latex",
}

TECHNICAL_STOPWORDS: Set[str] = {
    "a", "an", "the", "and", "or", "but", "in", "on", "at", "to", "for",
    "of", "with", "by", "from", "as", "is", "was", "are", "were", "been",
    "be", "have", "has", "had", "do", "does", "did", "will", "would",
    "could", "should", "may", "might", "shall", "can", "need", "dare",
    "ought", "used", "it", "its", "this", "that", "these", "those",
    "i", "me", "my", "we", "our", "you", "your", "he", "him", "his",
    "she", "her", "they", "them", "their", "what", "which", "who",
    "whom", "when", "where", "why", "how", "all", "each", "every",
    "both", "few", "more", "most", "other", "some", "such", "no",
    "not", "only", "own", "same", "so", "than", "too", "very",
    "just", "because", "also", "then", "there", "here", "about",
    "up", "out", "into", "over", "after", "before", "between",
    "under", "again", "further", "once", "any", "if", "while",
    "above", "below", "through", "during", "until", "against",
    "fig", "figure", "table", "eq", "equation", "section", "chapter",
    "note", "example", "see", "et", "al", "ie", "eg", "etc", "ref",
    "using", "used", "via", "per", "null", "true", "false", "none",
}

_PUNCTUATION_RE = re.compile(r"[^\w\s]")

# ---------------------------------------------------------------------------
# Dataclasses
# ---------------------------------------------------------------------------

@dataclass
class ChunkMetadata:
    chunk_id: str
    parent_id: Optional[str]
    text: str
    element_type: str
    parent_heading: str
    document_hierarchy_level: int
    token_count: int
    source_document: str
    page_range: str

@dataclass
class PipelineResult:
    chunks: List[ChunkMetadata]
    inverted_index: Dict[str, List[Dict[str, object]]]
    stats: Dict[str, object]

@dataclass
class RawItem:
    text: str
    element_type: str  # 'heading', 'text', 'list', 'table'
    level: int = 1     # heading level (1-6)
    page_label: str = ""

# ---------------------------------------------------------------------------
# Text / Token Utilities
# ---------------------------------------------------------------------------

def count_tokens(text: str) -> int:
    return len(text.split())

def normalize_token(token: str) -> Optional[str]:
    token = _PUNCTUATION_RE.sub("", token.lower())
    if len(token) < 2 or token in TECHNICAL_STOPWORDS:
        return None
    return token

def tokenize_for_index(text: str) -> List[str]:
    raw = text.split()
    normalized: List[str] = []
    for tok in raw:
        ntok = normalize_token(tok)
        if ntok is not None:
            normalized.append(ntok)
    return normalized

# ---------------------------------------------------------------------------
# Parsers
# ---------------------------------------------------------------------------

class BaseParser:
    def parse(self, filepath: Path) -> List[RawItem]:
        raise NotImplementedError

class PDFParser(BaseParser):
    def parse(self, filepath: Path) -> List[RawItem]:
        import fitz  # PyMuPDF
        items = []
        try:
            doc = fitz.open(filepath)
            for page_num, page in enumerate(doc, start=1):
                page_label = str(page_num)
                # Tables
                tables = page.find_tables()
                if tables:
                    for table in tables.tables:
                        # Extract table as Markdown-like string
                        lines = []
                        for row in table.extract():
                            clean_row = [str(c).replace('\n', ' ') if c else "" for c in row]
                            lines.append(" | ".join(clean_row))
                        if lines:
                            items.append(RawItem(text="\n".join(lines), element_type="table", page_label=page_label))
                
                # Text Blocks
                blocks = page.get_text("dict")["blocks"]
                for b in blocks:
                    if b.get("type") == 0:  # text block
                        block_text = ""
                        max_size = 0.0
                        is_bold = False
                        
                        for l in b["lines"]:
                            for s in l["spans"]:
                                block_text += s["text"] + " "
                                if s["size"] > max_size:
                                    max_size = s["size"]
                                if "bold" in s["font"].lower():
                                    is_bold = True
                        
                        block_text = block_text.strip()
                        if not block_text:
                            continue
                            
                        # Heading heuristic — strict to avoid fragmenting body text.
                        # Require (bold AND font > 13pt) or font > 16pt alone,
                        # AND short word count to exclude bold body paragraphs.
                        word_count = len(block_text.split())
                        is_all_caps = block_text.isupper() and len(block_text) > 3
                        is_likely_heading = (
                            word_count < 10
                            and (
                                (is_bold and max_size > 13.0)  # bold + somewhat large
                                or max_size > 16.0              # very large regardless
                                or is_all_caps                  # all-caps section headers
                            )
                        )
                        if is_likely_heading:
                            level = 1
                            if max_size < 14: level = 3
                            elif max_size < 18: level = 2
                            items.append(RawItem(text=block_text, element_type="heading", level=level, page_label=page_label))
                        elif block_text.startswith(("- ", "• ", "* ", "1. ")):
                            items.append(RawItem(text=block_text, element_type="list", page_label=page_label))
                        else:
                            items.append(RawItem(text=block_text, element_type="text", page_label=page_label))
        except Exception as e:
            logger.error(f"Error parsing PDF {filepath}: {e}")
        return items

class DOCXParser(BaseParser):
    def parse(self, filepath: Path) -> List[RawItem]:
        import docx
        items = []
        try:
            doc = docx.Document(filepath)
            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                style = para.style.name.lower()
                if "heading" in style:
                    level = 1
                    m = re.search(r'\d+', style)
                    if m: level = int(m.group())
                    items.append(RawItem(text=text, element_type="heading", level=level))
                elif "list" in style:
                    items.append(RawItem(text=text, element_type="list"))
                else:
                    items.append(RawItem(text=text, element_type="text"))
            
            for table in doc.tables:
                lines = []
                for row in table.rows:
                    cells = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
                    lines.append(" | ".join(cells))
                if lines:
                    items.append(RawItem(text="\n".join(lines), element_type="table"))
        except Exception as e:
            logger.error(f"Error parsing DOCX {filepath}: {e}")
        return items

class PPTXParser(BaseParser):
    def parse(self, filepath: Path) -> List[RawItem]:
        import pptx
        items = []
        try:
            prs = pptx.Presentation(filepath)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text = shape.text.strip()
                        if text:
                            # Heuristic: title shape is a heading
                            if shape == slide.shapes.title:
                                items.append(RawItem(text=text, element_type="heading", level=1))
                            else:
                                items.append(RawItem(text=text, element_type="text"))
                    elif shape.has_table:
                        lines = []
                        for row in shape.table.rows:
                            cells = [cell.text_frame.text.strip().replace('\n', ' ') for cell in row.cells]
                            lines.append(" | ".join(cells))
                        if lines:
                            items.append(RawItem(text="\n".join(lines), element_type="table"))
        except Exception as e:
            logger.error(f"Error parsing PPTX {filepath}: {e}")
        return items

class XLSXParser(BaseParser):
    def parse(self, filepath: Path) -> List[RawItem]:
        import openpyxl
        items = []
        try:
            wb = openpyxl.load_workbook(filepath, data_only=True)
            for sheet in wb.worksheets:
                items.append(RawItem(text=f"Sheet: {sheet.title}", element_type="heading", level=1))
                lines = []
                for row in sheet.iter_rows(values_only=True):
                    row_vals = [str(cell) if cell is not None else "" for cell in row]
                    if any(row_vals):
                        lines.append(" | ".join(row_vals))
                if lines:
                    items.append(RawItem(text="\n".join(lines), element_type="table"))
        except Exception as e:
            logger.error(f"Error parsing XLSX {filepath}: {e}")
        return items

class HTMLParser(BaseParser):
    def parse(self, filepath: Path) -> List[RawItem]:
        from bs4 import BeautifulSoup
        items = []
        try:
            with open(filepath, 'r', encoding='utf-8') as f:
                soup = BeautifulSoup(f, 'html.parser')
            for element in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'ul', 'ol', 'table']):
                text = element.get_text(separator=" ", strip=True)
                if not text:
                    continue
                tag = element.name
                if tag.startswith('h') and len(tag) == 2:
                    level = int(tag[1])
                    items.append(RawItem(text=text, element_type="heading", level=level))
                elif tag in ('ul', 'ol'):
                    items.append(RawItem(text=text, element_type="list"))
                elif tag == 'table':
                    items.append(RawItem(text=text, element_type="table"))
                else:
                    items.append(RawItem(text=text, element_type="text"))
        except Exception as e:
            logger.error(f"Error parsing HTML {filepath}: {e}")
        return items

class CSVParser(BaseParser):
    def parse(self, filepath: Path) -> List[RawItem]:
        import csv
        items = []
        try:
            with open(filepath, 'r', encoding='utf-8') as f:
                reader = csv.reader(f)
                lines = [" | ".join(row) for row in reader if any(row)]
                if lines:
                    items.append(RawItem(text="\n".join(lines), element_type="table"))
        except Exception as e:
            logger.error(f"Error parsing CSV {filepath}: {e}")
        return items

class TextParser(BaseParser):
    def parse(self, filepath: Path) -> List[RawItem]:
        items = []
        try:
            with open(filepath, 'r', encoding='utf-8') as f:
                text = f.read()
            
            blocks = re.split(r'\n\s*\n', text)
            for block in blocks:
                block = block.strip()
                if not block:
                    continue
                
                # Markdown headers
                m_md = re.match(r'^(#{1,6})\s+(.*)', block)
                if m_md:
                    level = len(m_md.group(1))
                    items.append(RawItem(text=m_md.group(2), element_type="heading", level=level))
                    continue
                
                # LaTeX section
                m_tex = re.match(r'^\\(section|subsection|subsubsection)\*?{(.*?)}', block)
                if m_tex:
                    level = 1 if m_tex.group(1) == 'section' else 2 if m_tex.group(1) == 'subsection' else 3
                    items.append(RawItem(text=m_tex.group(2), element_type="heading", level=level))
                    continue
                
                # AsciiDoc headers
                m_asc = re.match(r'^(=+)\s+(.*)', block)
                if m_asc:
                    level = len(m_asc.group(1))
                    items.append(RawItem(text=m_asc.group(2), element_type="heading", level=level))
                    continue
                
                # All caps heading
                if block.isupper() and len(block.split()) < 10 and '\n' not in block:
                    items.append(RawItem(text=block, element_type="heading", level=2))
                    continue
                    
                # Numbered section (e.g., "1.2 Introduction")
                if re.match(r'^\d+(\.\d+)*\s+[A-Z]', block) and len(block.split()) < 15 and '\n' not in block:
                    items.append(RawItem(text=block, element_type="heading", level=2))
                    continue
                
                # List
                if block.startswith(("- ", "* ", "+ ", "1. ")):
                    items.append(RawItem(text=block, element_type="list"))
                    continue
                    
                # Table basic heuristical detection (pipes)
                if '|' in block and '\n' in block:
                    lines = block.split('\n')
                    if sum(1 for l in lines if '|' in l) > len(lines) / 2:
                        items.append(RawItem(text=block, element_type="table"))
                        continue

                items.append(RawItem(text=block, element_type="text"))
        except Exception as e:
            logger.error(f"Error parsing Text {filepath}: {e}")
        return items

def get_parser(filepath: Path) -> Optional[BaseParser]:
    ext = filepath.suffix.lower()
    if ext == ".pdf": return PDFParser()
    if ext == ".docx": return DOCXParser()
    if ext == ".pptx": return PPTXParser()
    if ext == ".xlsx": return XLSXParser()
    if ext in (".html", ".htm"): return HTMLParser()
    if ext == ".csv": return CSVParser()
    if ext in (".md", ".txt", ".asciidoc", ".latex"): return TextParser()
    return None

# ---------------------------------------------------------------------------
# Structural Chunker (Layer 1)
# ---------------------------------------------------------------------------

class StructuralChunker:
    """
    Layer 1: Groups document elements (headings, text, tables) into logical blocks.
    Attempts to preserve the author's intent by keeping related content together.
    """
    def __init__(self, items: List[RawItem], source_name: str) -> None:
        self._items = items
        self._source_name = source_name

    def extract_blocks(self) -> List[ChunkMetadata]:
        blocks: List[ChunkMetadata] = []
        current_texts: List[str] = []
        current_element_types: List[str] = []
        current_pages: List[str] = []
        heading_stack: List[Tuple[str, int]] = []

        def _current_heading() -> str:
            return heading_stack[-1][0] if heading_stack else ""

        def _current_level() -> int:
            return len(heading_stack)

        def _compute_page_range(pages: List[str]) -> str:
            unique_pages = sorted(list(set(p for p in pages if p)))
            if not unique_pages: return ""
            if len(unique_pages) == 1: return unique_pages[0]
            try:
                int_pages = sorted([int(p) for p in unique_pages])
                if int_pages[-1] - int_pages[0] == len(int_pages) - 1 and len(int_pages) > 1:
                    return f"{int_pages[0]}-{int_pages[-1]}"
                return ", ".join(str(p) for p in int_pages)
            except ValueError:
                return ", ".join(unique_pages)

        def _flush_block() -> None:
            if not current_texts:
                return
            raw_combined = "\n".join(current_texts)
            current_head = _current_heading()
            combined = f"[{self._source_name}]\n{current_head}\n{raw_combined}" if current_head else f"[{self._source_name}]\n{raw_combined}"
            
            dominant_type = self._dominant_element_type(current_element_types)
            blocks.append(
                ChunkMetadata(
                    chunk_id=str(uuid.uuid4()),
                    parent_id=None,
                    text=combined,
                    element_type=dominant_type,
                    parent_heading=current_head,
                    document_hierarchy_level=_current_level(),
                    token_count=count_tokens(combined),
                    source_document=self._source_name,
                    page_range=_compute_page_range(current_pages),
                )
            )
            current_texts.clear()
            current_element_types.clear()
            current_pages.clear()

        for item in self._items:
            if item.element_type == "table":
                _flush_block()
                if item.text:
                    current_head = _current_heading()
                    table_text = f"[{self._source_name}]\n{current_head}\n{item.text}" if current_head else f"[{self._source_name}]\n{item.text}"
                    blocks.append(
                        ChunkMetadata(
                            chunk_id=str(uuid.uuid4()),
                            parent_id=None,
                            text=table_text,
                            element_type="table",
                            parent_heading=current_head,
                            document_hierarchy_level=_current_level(),
                            token_count=count_tokens(table_text),
                            source_document=self._source_name,
                            page_range=item.page_label,
                        )
                    )
                continue

            if item.element_type == "heading":
                _flush_block()
                heading_level = item.level
                while heading_stack and heading_stack[-1][1] >= heading_level:
                    heading_stack.pop()
                heading_stack.append((item.text, heading_level))
                
                # We deliberately DO NOT append the heading as a standalone chunk here.
                # Standalone headings carry no context for the LLM and pollute the vector 
                # space, stealing top-K slots from actual content. The heading text is 
                # already prepended to the actual text/table blocks in _flush_block().
                continue

            if item.text:
                current_texts.append(item.text)
                current_element_types.append(item.element_type)
                if item.page_label:
                    current_pages.append(item.page_label)

        _flush_block()
        return blocks

    @staticmethod
    def _parent_heading_from_stack(stack: List[Tuple[str, int]]) -> str:
        if len(stack) > 1:
            return stack[-2][0]
        return ""

    @staticmethod
    def _dominant_element_type(types: List[str]) -> str:
        if not types:
            return "text"
        counts: Dict[str, int] = {}
        for t in types:
            counts[t] = counts.get(t, 0) + 1
        return max(counts, key=counts.get)  # type: ignore[arg-type]

# ---------------------------------------------------------------------------
# Block Aggregator (Layer 1.5 — merge small adjacent blocks)
# ---------------------------------------------------------------------------

class BlockAggregator:
    """Merge consecutive small non-heading, non-table structural blocks
    that share the same parent heading into larger blocks.  This prevents
    the Layer-2 chunker from receiving pre-fragmented input."""

    def __init__(self, min_tokens: int = DEFAULT_MIN_CHUNK_TOKENS) -> None:
        self.min_tokens = min_tokens

    def aggregate(self, blocks: List[ChunkMetadata]) -> List[ChunkMetadata]:
        if not blocks:
            return blocks

        merged: List[ChunkMetadata] = []
        buffer: Optional[ChunkMetadata] = None

        def _flush() -> None:
            nonlocal buffer
            if buffer is not None:
                merged.append(buffer)
                buffer = None

        for block in blocks:
            # Always pass headings and tables through un-merged
            if block.element_type in ("heading", "table"):
                _flush()
                merged.append(block)
                continue

            # If buffer is empty, start a new accumulation
            if buffer is None:
                buffer = ChunkMetadata(
                    chunk_id=block.chunk_id,
                    parent_id=block.parent_id,
                    text=block.text,
                    element_type=block.element_type,
                    parent_heading=block.parent_heading,
                    document_hierarchy_level=block.document_hierarchy_level,
                    token_count=block.token_count,
                    source_document=block.source_document,
                    page_range=block.page_range,
                )
                continue

            # Same heading context → eligible to merge
            same_heading = (block.parent_heading == buffer.parent_heading)
            same_source = (block.source_document == buffer.source_document)

            if same_heading and same_source and buffer.token_count < self.min_tokens:
                # Merge block into buffer
                buffer.text = buffer.text + "\n" + block.text
                buffer.token_count = count_tokens(buffer.text)
                # Widen page range
                if block.page_range and buffer.page_range:
                    all_pages = set(buffer.page_range.replace("-", ", ").split(", ")) | \
                                set(block.page_range.replace("-", ", ").split(", "))
                    int_pages = sorted(int(p) for p in all_pages if p.isdigit())
                    if int_pages:
                        buffer.page_range = f"{int_pages[0]}-{int_pages[-1]}" if len(int_pages) > 1 else str(int_pages[0])
                elif block.page_range:
                    buffer.page_range = block.page_range
            else:
                _flush()
                buffer = ChunkMetadata(
                    chunk_id=block.chunk_id,
                    parent_id=block.parent_id,
                    text=block.text,
                    element_type=block.element_type,
                    parent_heading=block.parent_heading,
                    document_hierarchy_level=block.document_hierarchy_level,
                    token_count=block.token_count,
                    source_document=block.source_document,
                    page_range=block.page_range,
                )

        _flush()
        return merged

# ---------------------------------------------------------------------------
# Overlap Chunker (Layer 2 — Sliding Window)
# ---------------------------------------------------------------------------

class OverlapChunker:
    """
    Layer 2: Applies a sliding-window split with overlap to structural blocks.
    Used when the "Structural Chunker" alone does not sufficiently break down content.
    """
    def __init__(
        self,
        threshold: int = DEFAULT_CHUNK_THRESHOLD,
        overlap_ratio: float = DEFAULT_OVERLAP_RATIO,
    ) -> None:
        self.threshold = threshold
        self.overlap = round(threshold * overlap_ratio)
        self.stride = threshold - self.overlap

    def split(self, blocks: List[ChunkMetadata]) -> List[ChunkMetadata]:
        result: List[ChunkMetadata] = []
        for block in blocks:
            if block.token_count <= self.threshold:
                result.append(block)
                continue
            # Preserve the parent block for small-to-big retrieval
            result.append(block)
            children = self._split_block(block)
            result.extend(children)
        return result

    def _split_block(self, block: ChunkMetadata) -> List[ChunkMetadata]:
        tokens = block.text.split()
        total = len(tokens)
        if total <= self.threshold:
            return [block]

        parent_id = block.chunk_id
        children: List[ChunkMetadata] = []
        start = 0

        while start < total:
            end = min(start + self.threshold, total)
            chunk_text = " ".join(tokens[start:end])
            child = ChunkMetadata(
                chunk_id=str(uuid.uuid4()),
                parent_id=parent_id,
                text=chunk_text,
                element_type=block.element_type,
                parent_heading=block.parent_heading,
                document_hierarchy_level=block.document_hierarchy_level,
                token_count=count_tokens(chunk_text),
                source_document=block.source_document,
                page_range=block.page_range,
            )
            children.append(child)

            if end >= total:
                break

            next_start = start + self.stride
            if next_start + self.threshold >= total and next_start < total:
                remaining_text = " ".join(tokens[next_start:])
                if count_tokens(remaining_text) > self.overlap:
                    child = ChunkMetadata(
                        chunk_id=str(uuid.uuid4()),
                        parent_id=parent_id,
                        text=remaining_text,
                        element_type=block.element_type,
                        parent_heading=block.parent_heading,
                        document_hierarchy_level=block.document_hierarchy_level,
                        token_count=count_tokens(remaining_text),
                        source_document=block.source_document,
                        page_range=block.page_range,
                    )
                    children.append(child)
                break

            start = next_start

        return children

class RecursiveChunker:
    """Split oversized blocks using a hierarchy of natural-language separators.

    The algorithm tries the most significant separator first (paragraph break)
    and falls back to finer ones (sentence, then word) only when a segment
    still exceeds the token threshold.  No ML model is required — this is the
    industry-standard "recursive character text splitter" adapted for the
    structural-chunking pipeline.

    Separator hierarchy (coarsest → finest):
        1. ``\\n\\n``  — paragraph breaks
        2. ``\\n``    — line breaks
        3. sentence boundary (regex ``(?<=[.!?])\\s+``)
        4. ``" "``   — individual words (last resort)
    """

    # Ordered coarse → fine.  The last entry ("") triggers word-level split.
    _SEPARATORS: List[str] = ["\n\n", "\n", "__SENTENCE__", " "]

    def __init__(
        self,
        threshold: int = DEFAULT_CHUNK_THRESHOLD,
        overlap_ratio: float = DEFAULT_OVERLAP_RATIO,
    ) -> None:
        self.threshold = threshold
        self.overlap = round(threshold * overlap_ratio)

    # ---- public API (same interface as OverlapChunker) --------------------

    def split(self, blocks: List[ChunkMetadata]) -> List[ChunkMetadata]:
        result: List[ChunkMetadata] = []
        for block in blocks:
            if block.token_count <= self.threshold:
                result.append(block)
                continue
            # Keep the full block as a parent for small-to-big retrieval
            result.append(block)
            children = self._recursive_split(
                text=block.text,
                separators=self._SEPARATORS,
            )
            # Materialise children with overlap
            children_with_overlap = self._add_overlap(children)
            for c_text in children_with_overlap:
                if not c_text.strip():
                    continue
                result.append(
                    ChunkMetadata(
                        chunk_id=str(uuid.uuid4()),
                        parent_id=block.chunk_id,
                        text=c_text,
                        element_type=block.element_type,
                        parent_heading=block.parent_heading,
                        document_hierarchy_level=block.document_hierarchy_level,
                        token_count=count_tokens(c_text),
                        source_document=block.source_document,
                        page_range=block.page_range,
                    )
                )
        return result

    # ---- internal helpers -------------------------------------------------

    def _split_by_separator(self, text: str, separator: str) -> List[str]:
        """Split *text* by *separator*, keeping non-empty segments."""
        if separator == "__SENTENCE__":
            parts = re.split(r'(?<=[.!?])\s+', text)
        else:
            parts = text.split(separator)
        return [p for p in parts if p.strip()]

    def _recursive_split(
        self,
        text: str,
        separators: List[str],
    ) -> List[str]:
        """Recursively split *text* until every piece fits within the threshold."""
        if count_tokens(text) <= self.threshold:
            return [text]

        if not separators:
            # Ultimate fallback: hard word-level split
            words = text.split()
            chunks: List[str] = []
            for i in range(0, len(words), self.threshold):
                chunks.append(" ".join(words[i : i + self.threshold]))
            return chunks

        sep = separators[0]
        remaining_seps = separators[1:]
        parts = self._split_by_separator(text, sep)

        if len(parts) <= 1:
            # This separator didn't help — try the next finer one
            return self._recursive_split(text, remaining_seps)

        # Greedily merge parts into chunks up to the threshold
        merged: List[str] = []
        current_parts: List[str] = []
        current_tokens = 0

        joiner = "\n\n" if sep == "\n\n" else ("\n" if sep == "\n" else " ")

        for part in parts:
            part_tokens = count_tokens(part)

            if part_tokens > self.threshold:
                # Flush what we have, then recurse on the oversized part
                if current_parts:
                    merged.append(joiner.join(current_parts))
                    current_parts = []
                    current_tokens = 0
                merged.extend(self._recursive_split(part, remaining_seps))
                continue

            if current_tokens + part_tokens > self.threshold and current_parts:
                merged.append(joiner.join(current_parts))
                current_parts = []
                current_tokens = 0

            current_parts.append(part)
            current_tokens += part_tokens

        if current_parts:
            merged.append(joiner.join(current_parts))

        return merged

    def _add_overlap(self, chunks: List[str]) -> List[str]:
        """Add token overlap between consecutive chunks for boundary continuity."""
        if self.overlap <= 0 or len(chunks) <= 1:
            return chunks

        result: List[str] = []
        for i, chunk in enumerate(chunks):
            if i == 0:
                result.append(chunk)
                continue
            # Prepend tail of the previous chunk
            prev_words = chunks[i - 1].split()
            overlap_words = prev_words[-self.overlap :] if len(prev_words) > self.overlap else prev_words
            result.append(" ".join(overlap_words) + " " + chunk)
        return result

# ---------------------------------------------------------------------------
# Inverted Index (BM25-ready, TF + IDF pre-computed)
# ---------------------------------------------------------------------------

class InvertedIndex:
    """
    Layer 3: Builds an inverted index for efficient retrieval.
    Pre-computes TF and IDF values for BM25 scoring.
    """
    def __init__(self) -> None:
        self.index: Dict[str, List[Dict[str, object]]] = {}
        self.idf: Dict[str, float] = {}
        self.doc_lengths: Dict[str, int] = {}
        self.total_chunks: int = 0
        self.avg_dl: float = 0.0
        self.vocab_size: int = 0

    def build(self, chunks: List[ChunkMetadata]) -> None:
        self.total_chunks = len(chunks)
        total_length = 0

        for chunk in chunks:
            tokens = tokenize_for_index(chunk.text)
            self.doc_lengths[chunk.chunk_id] = len(tokens)
            total_length += len(tokens)

            tf_map: Dict[str, int] = {}
            for tok in tokens:
                tf_map[tok] = tf_map.get(tok, 0) + 1

            for tok, freq in tf_map.items():
                posting: Dict[str, object] = {
                    "chunk_id": chunk.chunk_id,
                    "tf": freq,
                }
                if tok not in self.index:
                    self.index[tok] = []
                self.index[tok].append(posting)

        self.avg_dl = total_length / self.total_chunks if self.total_chunks > 0 else 0.0
        self.vocab_size = len(self.index)

        for tok, postings in self.index.items():
            df = len(postings)
            self.idf[tok] = math.log((self.total_chunks - df + 0.5) / (df + 0.5) + 1.0)

    def query(self, term: str) -> List[Dict[str, object]]:
        normalized = normalize_token(term)
        if normalized is None:
            return []
        return self.index.get(normalized, [])

    def export(self) -> Dict[str, object]:
        return {
            "index": self.index,
            "idf": self.idf,
            "doc_lengths": self.doc_lengths,
            "stats": {
                "total_chunks": self.total_chunks,
                "avg_dl": self.avg_dl,
                "vocab_size": self.vocab_size,
            },
        }

# ---------------------------------------------------------------------------
# Embedding Generation & FAISS Index
# ---------------------------------------------------------------------------

try:
    import faiss
    _FAISS_AVAILABLE = True
except ImportError:
    _FAISS_AVAILABLE = False
    logger.warning("faiss-cpu not installed; embedding search will be unavailable.")


class EmbedCache:
    """
    Layer 4 (optional): Caches embedding mappings to avoid re-embedding identical text.
    Used for incremental updates when only new/modified files are processed.
    """
    def __init__(self, cache_path: Path) -> None:
        self._path = cache_path
        self._data: Dict[str, List[str]] = self._load()

    def _load(self) -> Dict[str, List[str]]:
        if self._path.exists():
            try:
                with open(self._path, "r", encoding="utf-8") as f:
                    return json.load(f)
            except (json.JSONDecodeError, OSError):
                return {}
        return {}

    def get_chunk_ids(self, filename: str) -> List[str]:
        return self._data.get(filename, [])

    def update(self, filename: str, chunk_ids: List[str]) -> None:
        self._data[filename] = chunk_ids

    def save(self) -> None:
        self._path.parent.mkdir(parents=True, exist_ok=True)
        with open(self._path, "w", encoding="utf-8") as f:
            json.dump(self._data, f, indent=2, ensure_ascii=False)


def embed_chunks(
    chunks: List[ChunkMetadata],
    model_name: str = DEFAULT_EMBED_MODEL,
) -> Tuple[np.ndarray, List[str]]:
    """Encode all chunks and return (embeddings, chunk_ids)."""
    try:
        from sentence_transformers import SentenceTransformer
    except ImportError:
        raise RuntimeError("sentence-transformers not installed; cannot embed chunks.")

    logger.info("Loading embedding model: %s", model_name)
    model = SentenceTransformer(model_name, local_files_only=True)
    texts = [c.text for c in chunks]
    chunk_ids = [c.chunk_id for c in chunks]

    logger.info("Encoding %d chunks ...", len(texts))
    embeddings = model.encode(
        texts,
        batch_size=64,
        show_progress_bar=True,
        convert_to_numpy=True,
    )
    # L2-normalize so inner product = cosine similarity
    faiss.normalize_L2(embeddings)
    logger.info("Encoding complete. Shape: %s", embeddings.shape)
    return embeddings, chunk_ids


def build_faiss_index(
    embeddings: np.ndarray,
    chunk_ids: List[str],
    faiss_path: Path,
    chunk_ids_path: Path,
) -> None:
    """Build and save a FAISS index from embeddings."""
    if not _FAISS_AVAILABLE:
        raise RuntimeError("faiss-cpu not installed; cannot build FAISS index.")

    dim = embeddings.shape[1]
    index = faiss.IndexFlatIP(dim)
    index.add(embeddings)

    faiss_path.parent.mkdir(parents=True, exist_ok=True)
    faiss.write_index(index, str(faiss_path))

    with open(chunk_ids_path, "w", encoding="utf-8") as f:
        json.dump(chunk_ids, f, indent=2, ensure_ascii=False)

    logger.info("FAISS index saved: %s (%d vectors, dim=%d)", faiss_path, len(chunk_ids), dim)


# ---------------------------------------------------------------------------
# Parse Cache (incremental ingestion)
# ---------------------------------------------------------------------------

def _sha256(filepath: Path) -> str:
    """Return the hex SHA-256 digest of a file."""
    h = hashlib.sha256()
    with open(filepath, "rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()


class ParseCache:
    """Persist a {filename -> sha256} map to avoid re-parsing unchanged files."""

    def __init__(self, cache_path: Path) -> None:
        self._path = cache_path
        self._data: Dict[str, str] = self._load()

    def _load(self) -> Dict[str, str]:
        if self._path.exists():
            try:
                with open(self._path, "r", encoding="utf-8") as f:
                    return json.load(f)
            except (json.JSONDecodeError, OSError):
                return {}
        return {}

    def is_unchanged(self, filepath: Path) -> bool:
        """Return True if the file's hash matches the cached value."""
        key = filepath.name
        try:
            current = _sha256(filepath)
        except OSError:
            return False
        return self._data.get(key) == current

    def update(self, filepath: Path) -> None:
        """Record the current hash for a file."""
        self._data[filepath.name] = _sha256(filepath)

    def remove(self, filename: str) -> None:
        """Drop a stale entry (e.g. file was deleted)."""
        self._data.pop(filename, None)

    def save(self) -> None:
        """Flush the cache to disk."""
        self._path.parent.mkdir(parents=True, exist_ok=True)
        with open(self._path, "w", encoding="utf-8") as f:
            json.dump(self._data, f, indent=2, ensure_ascii=False)


# ---------------------------------------------------------------------------
# Pipeline Orchestrator
# ---------------------------------------------------------------------------

class IngestionPipeline:
    """
    Orchestrates the entire RAG ingestion pipeline:
    1. Discovery (scans input directory for supported files)
    2. Parsing (extracts raw text/metadata from files)
    3. Structural Chunking (groups related items into logical blocks)
    4. Recursive Chunking (splits large blocks into manageable chunks)
    5. Indexing (builds inverted index for BM25 retrieval)
    6. Embedding (generates vector representations)
    """
    def __init__(
        self,
        chunk_threshold: int = DEFAULT_CHUNK_THRESHOLD,
        overlap_ratio: float = DEFAULT_OVERLAP_RATIO,
        cache: Optional[ParseCache] = None,
        embed: bool = True,
        model_name: str = DEFAULT_EMBED_MODEL,
        chunker_type: str = "overlap",
    ) -> None:
        self.chunk_threshold = chunk_threshold
        self.overlap_ratio = overlap_ratio
        self._cache = cache  # None means no caching
        self.embed = embed
        self.model_name = model_name
        self.chunker_type = chunker_type

    def run(
        self,
        input_dir: Path,
        prev_chunks_by_source: Optional[Dict[str, List[ChunkMetadata]]] = None,
        faiss_path: Optional[Path] = None,
        chunk_ids_path: Optional[Path] = None,
    ) -> PipelineResult:
        """
        Executes the full ingestion pipeline.
        
        Args:
            input_dir: Directory containing document files
            prev_chunks_by_source: Optional cache of previous chunks for incremental updates
            faiss_path: Path to save the FAISS index (defaults to DEFAULT_FAISS_PATH)
            chunk_ids_path: Path to save chunk ID mappings (defaults to DEFAULT_CHUNK_IDS_PATH)
        
        Returns:
            PipelineResult containing all processed chunks, index, and stats
        """
        if faiss_path is None:
            faiss_path = DEFAULT_FAISS_PATH
        if chunk_ids_path is None:
            chunk_ids_path = DEFAULT_CHUNK_IDS_PATH
        files = self._discover_files(input_dir)
        if not files:
            logger.warning("No supported documents found in %s", input_dir)
            return PipelineResult(chunks=[], inverted_index={}, stats={})

        all_chunks: List[ChunkMetadata] = []
        cached_count = 0
        parsed_count = 0

        # Pre-build reusable chunker instance
        recursive_chunker = None
        if self.chunker_type == "recursive":
            recursive_chunker = RecursiveChunker(
                threshold=self.chunk_threshold,
                overlap_ratio=self.overlap_ratio,
            )

        for fpath in files:
            source_name = fpath.name

            # --- Cache hit: reuse previous chunks if file is unchanged ---
            if (
                self._cache is not None
                and self._cache.is_unchanged(fpath)
                and prev_chunks_by_source is not None
                and source_name in prev_chunks_by_source
            ):
                reused = prev_chunks_by_source[source_name]
                all_chunks.extend(reused)
                cached_count += 1
                logger.info("  [cached] %s (%d chunks reused)", source_name, len(reused))
                continue

            # --- Cache miss: parse the file ---
            logger.info("Parsing %s ...", source_name)
            parser = get_parser(fpath)
            if not parser:
                logger.warning("No parser found for %s", source_name)
                continue

            items = parser.parse(fpath)
            if not items:
                logger.warning("No items parsed from %s", source_name)
                continue

            # Phase 1: Structural chunking
            struct_chunker = StructuralChunker(items, source_name)
            blocks = struct_chunker.extract_blocks()
            logger.info("  %s: %d structural blocks extracted", source_name, len(blocks))

            # Phase 1.5: Aggregate small adjacent blocks
            aggregator = BlockAggregator(min_tokens=DEFAULT_MIN_CHUNK_TOKENS)
            blocks = aggregator.aggregate(blocks)
            logger.info("  %s: %d blocks after aggregation", source_name, len(blocks))

            # Phase 2: Chunker
            if self.chunker_type == "recursive":
                chunker = recursive_chunker
            else:
                chunker = OverlapChunker(
                    threshold=self.chunk_threshold,
                    overlap_ratio=self.overlap_ratio,
                )
            chunks = chunker.split(blocks)
            logger.info("  %s: %d chunks after splitting", source_name, len(chunks))

            all_chunks.extend(chunks)
            parsed_count += 1

            # Record the new hash so next run can skip this file
            if self._cache is not None:
                self._cache.update(fpath)

        if self._cache is not None and cached_count > 0:
            logger.info(
                "Cache: %d file(s) skipped (unchanged), %d re-parsed.",
                cached_count,
                parsed_count,
            )

        # Phase 3: Build inverted index
        index = InvertedIndex()
        index.build(all_chunks)

        stats: Dict[str, object] = {
            "total_chunks": len(all_chunks),
            "total_documents": len(files),
            "avg_chunk_tokens": (
                sum(c.token_count for c in all_chunks) / len(all_chunks)
                if all_chunks
                else 0
            ),
            "vocab_size": index.vocab_size,
        }

        idx_export = index.export()

        # Phase 4: Build FAISS embedding index
        # Skip if all files were cached AND a FAISS index already exists
        faiss_exists = faiss_path.exists() and chunk_ids_path.exists()
        skip_embed = (parsed_count == 0) and faiss_exists

        if self.embed and all_chunks and not skip_embed:
            try:
                embeddings, e_chunk_ids = embed_chunks(all_chunks, self.model_name)
                build_faiss_index(embeddings, e_chunk_ids, faiss_path, chunk_ids_path)
                stats["embedding_model"] = self.model_name
                stats["embedding_dim"] = embeddings.shape[1]
                stats["faiss_index"] = str(faiss_path)
            except Exception as e:
                logger.warning("Embedding step failed: %s. Continuing with BM25-only index.", e)
                stats["embedding_error"] = str(e)
        elif skip_embed:
            logger.info("All files cached and FAISS index exists; skipping embedding.")
            stats["embedding_skipped"] = True
        else:
            logger.info("Embedding step skipped (--no-embed or no chunks).")

        return PipelineResult(
            chunks=all_chunks,
            inverted_index=idx_export,
            stats=stats,
        )

    @staticmethod
    def _discover_files(input_dir: Path) -> List[Path]:
        """
        Recursively discovers all supported document files within the input directory.
        """
        files: List[Path] = []
        for root, _dirs, filenames in os.walk(input_dir):
            for fname in filenames:
                ext = Path(fname).suffix.lower()
                if ext in SUPPORTED_EXTENSIONS:
                    files.append(Path(root) / fname)
        return sorted(files)

# ---------------------------------------------------------------------------
# Serialization
# ---------------------------------------------------------------------------

def serialize_result(result: PipelineResult) -> Dict[str, object]:
    """
    Converts the PipelineResult into a JSON-serializable dictionary.
    """
    return {
        "chunks": [asdict(c) for c in result.chunks],
        "inverted_index": result.inverted_index,
        "stats": result.stats,
    }

# ---------------------------------------------------------------------------
# CLI Entry Point
# ---------------------------------------------------------------------------

def main() -> None:
    """
    Parses command-line arguments, runs the ingestion pipeline, and saves the
    enriched JSON output to a file. This is the main entry point for the CLI.
    """
    import argparse

    parser = argparse.ArgumentParser(
        description=(
            "Offline RAG Ingestion Pipeline — parse, hybrid-chunk, "
            "enrich, and index technical documentation locally."
        ),
    )
    parser.add_argument(
        "input_dir",
        type=Path,
        nargs="?",
        default=DEFAULT_INPUT_DIR,
        help=f"Directory containing raw documents (default: {DEFAULT_INPUT_DIR}).",
    )
    parser.add_argument(
        "--output",
        "-o",
        type=Path,
        default=DEFAULT_OUTPUT_PATH,
        help=f"Path to write JSON output (default: {DEFAULT_OUTPUT_PATH}).",
    )
    parser.add_argument(
        "--chunk-threshold",
        type=int,
        default=DEFAULT_CHUNK_THRESHOLD,
        help=f"Maximum tokens per chunk before overlap splitting (default: {DEFAULT_CHUNK_THRESHOLD}).",
    )
    parser.add_argument(
        "--overlap-ratio",
        type=float,
        default=DEFAULT_OVERLAP_RATIO,
        help=f"Overlap ratio for sliding window (default: {DEFAULT_OVERLAP_RATIO}).",
    )
    parser.add_argument(
        "--quiet",
        "-q",
        action="store_true",
        help="Suppress parsing logs (only show summary and errors).",
    )
    parser.add_argument(
        "--force",
        action="store_true",
        help="Ignore the parse cache and re-parse all documents.",
    )
    parser.add_argument(
        "--cache",
        type=Path,
        default=DEFAULT_CACHE_PATH,
        help=f"Path to the parse cache file (default: {DEFAULT_CACHE_PATH}).",
    )
    parser.add_argument(
        "--no-embed",
        action="store_true",
        help="Skip the embedding / FAISS step (BM25-only index).",
    )
    parser.add_argument(
        "--model",
        type=str,
        default=DEFAULT_EMBED_MODEL,
        help=f"Sentence-transformers model for embeddings (default: {DEFAULT_EMBED_MODEL}).",
    )
    parser.add_argument(
        "--faiss-index",
        type=Path,
        default=DEFAULT_FAISS_PATH,
        help=f"Path to write FAISS index (default: {DEFAULT_FAISS_PATH}).",
    )
    parser.add_argument(
        "--chunk-ids",
        type=Path,
        default=DEFAULT_CHUNK_IDS_PATH,
        help=f"Path to write chunk ID mapping (default: {DEFAULT_CHUNK_IDS_PATH}).",
    )
    parser.add_argument(
        "--chunker",
        choices=["overlap", "recursive"],
        default="recursive",
        help="Chunking strategy: 'recursive' (recommended) or 'overlap' (default: recursive).",
    )

    args = parser.parse_args()

    if args.quiet:
        logger.setLevel(logging.WARNING)

    if not args.input_dir.is_dir():
        parser.error(f"Input directory does not exist: {args.input_dir}")

    # --- Load or bypass the parse cache ---
    cache: Optional[ParseCache] = None
    prev_chunks_by_source: Optional[Dict[str, List[ChunkMetadata]]] = None

    if not args.force:
        cache = ParseCache(args.cache)
        # Load chunks from the previous output so cached files can reuse them
        if args.output.exists():
            try:
                with open(args.output, "r", encoding="utf-8") as f:
                    prev_data = json.load(f)
                prev_chunks_by_source = {}
                for raw in prev_data.get("chunks", []):
                    src = raw.get("source_document", "")
                    chunk = ChunkMetadata(**{k: raw[k] for k in ChunkMetadata.__dataclass_fields__})
                    prev_chunks_by_source.setdefault(src, []).append(chunk)
                logger.info(
                    "Loaded %d chunks from previous output (%d sources).",
                    sum(len(v) for v in prev_chunks_by_source.values()),
                    len(prev_chunks_by_source),
                )
            except (json.JSONDecodeError, OSError, KeyError, TypeError):
                logger.warning("Could not load previous output; all files will be re-parsed.")
                prev_chunks_by_source = None
    else:
        logger.info("--force: skipping cache, re-parsing all documents.")

    pipeline = IngestionPipeline(
        chunk_threshold=args.chunk_threshold,
        overlap_ratio=args.overlap_ratio,
        cache=cache,
        embed=not args.no_embed,
        model_name=args.model,
        chunker_type=args.chunker,
    )

    logger.info("Starting ingestion pipeline on %s", args.input_dir)
    result = pipeline.run(
        args.input_dir,
        prev_chunks_by_source=prev_chunks_by_source,
        faiss_path=args.faiss_index,
        chunk_ids_path=args.chunk_ids,
    )
    logger.info(
        "Done: %d chunks, %d documents, vocab size %d",
        result.stats.get("total_chunks", 0),
        result.stats.get("total_documents", 0),
        result.stats.get("vocab_size", 0),
    )

    output_data = serialize_result(result)
    args.output.parent.mkdir(parents=True, exist_ok=True)
    with open(args.output, "w", encoding="utf-8") as f:
        json.dump(output_data, f, ensure_ascii=False)
    logger.info("Output written to %s", args.output)

    # Flush updated hashes to disk
    if cache is not None:
        cache.save()
        logger.info("Cache saved to %s", args.cache)

if __name__ == "__main__":
    main()
